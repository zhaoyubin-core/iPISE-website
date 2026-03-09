import os
from pptx import Presentation

pptx_path = "网站(1).pptx"
output_dir = "extracted_assets"

if not os.path.exists(output_dir):
    os.makedirs(output_dir)

print(f"Loading {pptx_path}...")
try:
    prs = Presentation(pptx_path)
    text_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = f"--- Slide {i+1} ---\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
            
            # Extract images
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_filename = f"{output_dir}/slide_{i+1}_img_{shape.shape_id}.{image_ext}"
                    with open(image_filename, "wb") as f:
                        f.write(image_bytes)
                except Exception as img_e:
                    print(f"Warning: Could not extract image on slide {i+1}: {img_e}")
                    
        text_content.append(slide_text)
        print(f"Processed slide {i+1}/{len(prs.slides)}")
        
    with open(f"{output_dir}/content_extracted.txt", "w", encoding="utf-8") as f:
        f.write("\n".join(text_content))
        
    print(f"Extraction complete! Text and images saved to {output_dir}")
except Exception as e:
    print(f"Error extracting PPTX: {e}")
